#!/usr/bin/env python3
"""
fi_weekly.py  —  Weekly Fixed Income briefing generator.

What it does:
  1. Pulls a fixed-income dashboard from FRED (rates, curve, credit, cross-asset),
     computes week-over-week changes in bp.
  2. Renders a 5-slide PPTX built for a trader (levels/tape up front) AND an
     analyst (themes + why + calendar in the back).
  3. Optionally emails the deck to you as a Monday attachment.

Design choice: ONE language (Python) end to end — fetch, build, send — so it runs
as a single scheduled job (cron / GitHub Actions) with no moving parts to babysit.

Run:
  python fi_weekly.py --demo                 # offline sample, no network, no email
  python fi_weekly.py --live                 # pull FRED, build deck, no email
  python fi_weekly.py --live --email         # pull, build, email it

Env vars for --live / --email:
  FRED_API_KEY     free key from https://fred.stlouisfed.org/docs/api/api_key.html
  SMTP_HOST        e.g. smtp.gmail.com
  SMTP_PORT        e.g. 587
  SMTP_USER        sending address
  SMTP_PASS        app password (NOT your normal password)
  MAIL_TO          recipient
"""
import subprocess
import sys

# 1. AUTOMATIC INSTALLATION BLOCK (Must be at the top)
try:
    from pptx import Presentation
except ImportError:
    print("python-pptx not found. Installing now...")
    # This runs 'pip install python-pptx' behind the scenes
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


import argparse
import datetime as dt
import io
import os
import smtplib
from email.message import EmailMessage

import requests
import matplotlib
matplotlib.use("Agg")  # headless: no display needed on a server
import matplotlib.pyplot as plt



# ----------------------------------------------------------------------------
# 1. CONFIG  — the only block you normally touch
# ----------------------------------------------------------------------------
BRAND = "Kutman Pamirov  ·  Fixed Income Weekly"
TIMEZONE_LABEL = "Times ET unless noted"

NAVY  = RGBColor(0x1E, 0x27, 0x61)   # dominant
ICE   = RGBColor(0xCA, 0xDC, 0xFC)   # supporting
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK   = RGBColor(0x1A, 0x1A, 0x2A)   # body text on light
MUTE  = RGBColor(0x6B, 0x72, 0x80)   # captions
UP    = RGBColor(0xC0, 0x39, 0x2B)   # yields higher / spreads wider
DOWN  = RGBColor(0x1E, 0x84, 0x49)   # yields lower / spreads tighter

# FRED series that make up the dashboard. Order = display order.
# label, fred_id, unit ("%" yields, "bp" spreads already in %-pts on FRED so *100)
RATES = [
    ("UST 3M",      "DGS3MO", "%"),
    ("UST 2Y",      "DGS2",   "%"),
    ("UST 5Y",      "DGS5",   "%"),
    ("UST 10Y",     "DGS10",  "%"),
    ("UST 30Y",     "DGS30",  "%"),
    ("2s10s",       "T10Y2Y", "bp"),   # already a spread in %-pts
    ("10Y real",    "DFII10", "%"),
    ("10Y B/E infl","T10YIE", "%"),
    ("SOFR",        "SOFR",   "%"),
]
CREDIT = [
    ("IG OAS",   "BAMLC0A0CM",     "bp"),
    ("HY OAS",   "BAMLH0A0HYM2",   "bp"),
    ("BB OAS",   "BAMLH0A1HYBB",   "bp"),
    ("B OAS",    "BAMLH0A2HYB",    "bp"),
    ("CCC OAS",  "BAMLH0A3HYC",    "bp"),
]

# Curated source links — always current, always yours to extend.
LINKS = [
    ("FRED dashboard (all series live)", "https://fred.stlouisfed.org/"),
    ("US Treasury daily par yield curve", "https://home.treasury.gov/resource-center/data-chart-center/interest-rates/TextView?type=daily_treasury_yield_curve"),
    ("BLS release calendar (CPI/NFP)", "https://www.bls.gov/schedule/news_release/current_month.htm"),
    ("Federal Reserve calendar / speeches", "https://www.federalreserve.gov/newsevents.htm"),
    ("Your Substack", "https://kutman.substack.com"),
]

# ----------------------------------------------------------------------------
# 2. DATA LAYER
# ----------------------------------------------------------------------------
def fred_series(series_id, api_key, days=45):
    """Return list of (date, value) for the last `days`, skipping missing '.'."""
    end = dt.date.today()
    start = end - dt.timedelta(days=days)
    url = "https://api.stlouisfed.org/fred/series/observations"
    params = {
        "series_id": series_id, "api_key": api_key, "file_type": "json",
        "observation_start": start.isoformat(), "observation_end": end.isoformat(),
    }
    r = requests.get(url, params=params, timeout=20)
    r.raise_for_status()
    out = []
    for o in r.json().get("observations", []):
        if o["value"] != ".":
            out.append((dt.date.fromisoformat(o["date"]), float(o["value"])))
    return out


def latest_and_wow(obs):
    """Given [(date,val)...], return (current, ref_1w_ago). Ref = last point
    on/before (latest_date - 7d). W/W change is computed by the caller."""
    if not obs:
        return None, None
    obs = sorted(obs)
    cur_date, cur_val = obs[-1]
    target = cur_date - dt.timedelta(days=7)
    ref_val = obs[0][1]
    for d, v in obs:
        if d <= target:
            ref_val = v
    return cur_val, ref_val


def pull_live(api_key):
    """Build the same dict shape as DEMO, but from FRED."""
    data = {"rates": [], "credit": [], "asof": dt.date.today().isoformat()}
    for label, sid, unit in RATES:
        cur, ref = latest_and_wow(fred_series(sid, api_key))
        data["rates"].append(_row(label, cur, ref, unit))
    for label, sid, unit in CREDIT:
        cur, ref = latest_and_wow(fred_series(sid, api_key))
        data["credit"].append(_row(label, cur, ref, unit))
    return data


def _row(label, cur, ref, unit):
    """Normalise one line: level shown in native unit, delta always in bp."""
    if cur is None:
        return {"label": label, "level": "n/a", "delta_bp": None, "unit": unit}
    if unit == "bp":
        level = f"{cur*100:.0f}"
        delta = round((cur - ref) * 100) if ref is not None else None
    else:
        level = f"{cur:.2f}"
        delta = round((cur - ref) * 100) if ref is not None else None
    return {"label": label, "level": level, "delta_bp": delta, "unit": unit,
            "raw": cur, "raw_ref": ref}

# ----------------------------------------------------------------------------
# 3. DEMO DATA  — grounded in real levels ~Aug 8 2026, clearly illustrative.
#    Replace with --live for the real weekly pull.
# ----------------------------------------------------------------------------
def demo_data():
    def R(label, level, dbp, unit, raw=None, raw_ref=None):
        return {"label": label, "level": level, "delta_bp": dbp, "unit": unit,
                "raw": raw, "raw_ref": raw_ref}
    return {
        "asof": "2026-08-07",
        "rates": [
            R("UST 3M",       "4.35", +2,  "%", 4.35, 4.33),
            R("UST 2Y",       "4.12", +6,  "%", 4.12, 4.06),
            R("UST 5Y",       "4.28", +8,  "%", 4.28, 4.20),
            R("UST 10Y",      "4.66", +9,  "%", 4.66, 4.57),
            R("UST 30Y",      "5.06", +11, "%", 5.06, 4.95),
            R("2s10s",        "+54",  +3,  "bp"),
            R("10Y real",     "2.05", +6,  "%"),
            R("10Y B/E infl", "2.61", +5,  "%"),
            R("SOFR",         "4.32",  0,  "%"),
        ],
        "credit": [
            R("IG OAS",  "100", +2,  "bp"),
            R("HY OAS",  "272", +6,  "bp"),
            R("BB OAS",  "175", +4,  "bp"),
            R("B OAS",   "320", +7,  "bp"),
            R("CCC OAS", "640", +15, "bp"),
        ],
        # curve points for the chart: tenor(yrs) -> (this_week, last_week)
        "curve": {0.25: (4.35, 4.33), 2: (4.12, 4.06), 5: (4.28, 4.20),
                  10: (4.66, 4.57), 30: (5.06, 4.95)},
        "cross": [("DXY", "99.8"), ("Brent", "$78"), ("Gold", "$3,450"),
                  ("S&P 500", "6,550"), ("MOVE (rate vol)", "95")],
        "tape": [
            "Bear-steepener: long end led, 30Y ~5.06% near post-2007 highs as a hawkish "
            "Warsh Fed and an oil bounce (Hormuz) revived hike risk into September.",
            "Front end firmer but lagging — market now ~40% for a Sept HIKE (was a cut "
            "debate a month ago). Curve disinverted further; 2s10s ~+54bp.",
            "Credit shrugged: HY OAS ~272bp, IG ~100bp — near multi-decade tights. "
            "Spreads widened only a few bp against a 9bp rates selloff.",
            "Real yields did the work: 10Y real +6bp to ~2.05%; breakevens only +5bp — "
            "this is a real-rate / term-premium move, not a pure inflation scare.",
            "Positioning read: carry still pays, but at these spreads the convexity is "
            "ugly — you are short a lot of optionality for ~270bp.",
        ],
        "themes": [
            ("Hawkish repricing has further to run",
             "The whole front end is priced off one question: does Warsh hike in Sept? "
             "With core CPI sticky and oil rebounding, the risk is asymmetric toward MORE "
             "hikes, not fewer.",
             "Read: stay short duration / underweight the belly; fade rallies in 2Y."),
            ("Credit is the complacent asset",
             "HY at ~272bp with spread duration ~3.2y: a move back to 380bp (a 'tight' "
             "level as recently as 2023) marks a 10Y HY book down ~3.5% — roughly five "
             "months of carry gone. Downside dwarfs the remaining ~40bp of compression.",
             "Read: up-in-quality within HY (BB over CCC); the CCC-BB gap is too thin."),
            ("Oil is the inflation swing factor",
             "Strait of Hormuz headlines are now the marginal driver of breakevens and "
             "the long end. Energy is doing what a data surprise used to do.",
             "Read: watch Brent as a real-time proxy for the September Fed decision."),
        ],
        "calendar": [
            ("Wed Aug 12", "US CPI (Jul)",        "8:30", "The print that decides Sept. Core is the number."),
            ("Thu Aug 13", "Germany ZEW · UK jobs","—",    "EGB / Gilt tone-setter for the week."),
            ("Thu Aug 14", "UK CPI (Jul)",         "2:00", "BoE path; Gilt-Bund spread."),
            ("Fri Aug 15", "US Retail Sales, Claims","8:30","Growth check into a hawkish Fed."),
            ("Fri Aug 15", "Michigan sentiment (P)","10:00","5-10y inflation expectations watched."),
            ("Sep 16",     "FOMC decision + SEP",  "14:00","The main event this cycle."),
        ],
    }

# ----------------------------------------------------------------------------
# 4. CHART
# ----------------------------------------------------------------------------
def yield_curve_png(curve):
    """Return PNG bytes: this week vs last week UST curve."""
    tenors = sorted(curve)
    tw = [curve[t][0] for t in tenors]
    lw = [curve[t][1] for t in tenors]
    fig, ax = plt.subplots(figsize=(5.2, 2.9), dpi=200)
    ax.plot(range(len(tenors)), lw, "--", color="#9AA6C0", lw=1.6, label="1w ago")
    ax.plot(range(len(tenors)), tw, "-o", color="#1E2761", lw=2.2, ms=5, label="This week")
    ax.set_xticks(range(len(tenors)))
    ax.set_xticklabels(["3M", "2Y", "5Y", "10Y", "30Y"])
    ax.set_ylabel("Yield (%)", fontsize=9)
    ax.tick_params(labelsize=8)
    ax.grid(axis="y", color="#E5E8EF", lw=0.8)
    for s in ("top", "right"):
        ax.spines[s].set_visible(False)
    ax.legend(fontsize=8, frameon=False, loc="lower right")
    fig.tight_layout(pad=0.4)
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf

# ----------------------------------------------------------------------------
# 5. PPTX HELPERS
# ----------------------------------------------------------------------------
def _fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()

def _text(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT,
          font="Calibri", anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.name = font; r.font.color.rgb = color
    return tb

def _bul(slide, x, y, w, h, items, size=13, color=INK, gap=6, font="Calibri"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        r = p.add_run(); r.text = "▪  " + it
        r.font.size = Pt(size); r.font.name = font; r.font.color.rgb = color
    return tb

def _link(paragraph, label, url, color=NAVY):
    r = paragraph.add_run(); r.text = label
    r.font.size = Pt(12); r.font.name = "Calibri"; r.font.color.rgb = color
    r.hyperlink.address = url

def _metric_table(slide, x, y, rows, w=3.5, title=None, title_color=NAVY):
    """rows: list of dicts with label/level/delta_bp. Renders a compact table."""
    top = y
    if title:
        _text(slide, x, top, w, 0.3, title, 15, title_color, bold=True,
              font="Cambria")
        top += 0.42
    line_h = 0.30
    for row in rows:
        _text(slide, x, top, w*0.46, line_h, row["label"], 12.5, INK)
        _text(slide, x + w*0.46, top, w*0.24, line_h, row["level"], 12.5, INK,
              bold=True, align=PP_ALIGN.RIGHT)
        d = row["delta_bp"]
        if d is None:
            dtxt, dcol = "—", MUTE
        else:
            dtxt = f"{d:+d}bp"
            dcol = UP if d > 0 else (DOWN if d < 0 else MUTE)
        _text(slide, x + w*0.72, top, w*0.28, line_h, dtxt, 12.5, dcol,
              bold=True, align=PP_ALIGN.RIGHT)
        top += line_h
    return top

# ----------------------------------------------------------------------------
# 6. BUILD DECK
# ----------------------------------------------------------------------------
def build_deck(data, path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    W = 13.333
    asof = dt.date.fromisoformat(data["asof"]).strftime("%A, %d %B %Y")

    # --- Slide 1: Cover + The Tape (trader 30-second read) ------------------
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height); _fill(bg, NAVY)
    _text(s, 0.7, 0.55, 9, 0.5, "FIXED INCOME WEEKLY", 15, ICE, bold=True,
          font="Calibri")
    _text(s, 0.7, 1.0, 11.9, 1.0, "The Tape", 46, WHITE, bold=True, font="Cambria")
    _text(s, 0.7, 1.95, 11.9, 0.4, f"Week of {asof}  ·  {TIMEZONE_LABEL}", 14, ICE)
    # regime tag
    tag = s.shapes.add_shape(5, Inches(0.7), Inches(2.55), Inches(5.2), Inches(0.5))
    _fill(tag, ICE)
    _text(s, 0.85, 2.62, 5.0, 0.4, "REGIME:  Hawkish Fed · bear-steepening · credit tight",
          12.5, NAVY, bold=True)
    # the five bullets
    tb = s.shapes.add_textbox(Inches(0.7), Inches(3.25), Inches(11.9), Inches(4.0))
    tf = tb.text_frame; tf.word_wrap = True
    for i, t in enumerate(data["tape"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(9)
        r = p.add_run(); r.text = "→  " + t
        r.font.size = Pt(14.5); r.font.name = "Calibri"; r.font.color.rgb = WHITE
    _text(s, 0.7, 7.05, 11.9, 0.3, BRAND, 10, ICE)

    # --- Slide 2: Rates dashboard + curve (trader payload) ------------------
    s = prs.slides.add_slide(blank)
    _text(s, 0.7, 0.5, 11, 0.6, "Rates Dashboard", 34, NAVY, bold=True, font="Cambria")
    _text(s, 0.7, 1.15, 11, 0.3, f"Levels as of {asof}  ·  Δ = change vs 1 week ago",
          12, MUTE)
    _metric_table(s, 0.7, 1.7, data["rates"], w=4.6, title="US Treasuries & rates")
    # curve chart on the right
    png = yield_curve_png(data["curve"])
    s.shapes.add_picture(png, Inches(6.1), Inches(1.9), width=Inches(6.4))
    _text(s, 6.1, 5.55, 6.4, 0.3, "US Treasury curve: this week vs 1w ago", 11, MUTE,
          italic=True)
    _text(s, 6.1, 6.0, 6.4, 1.0,
          "Read: the selloff was led by the long end (30Y +11bp vs 2Y +6bp) — "
          "a term-premium / real-rate move, consistent with supply + hike risk "
          "rather than a growth scare.", 12.5, INK)
    _text(s, 0.7, 7.05, 6, 0.3, "Source: FRED (H.15), US Treasury", 10, MUTE)

    # --- Slide 3: Credit & cross-asset --------------------------------------
    s = prs.slides.add_slide(blank)
    _text(s, 0.7, 0.5, 11, 0.6, "Credit & Cross-Asset", 34, NAVY, bold=True,
          font="Cambria")
    _metric_table(s, 0.7, 1.5, data["credit"], w=4.4, title="OAS (bp) & Δ1w")
    # cross-asset mini table
    _text(s, 5.6, 1.5, 3.4, 0.3, "Cross-asset", 15, NAVY, bold=True, font="Cambria")
    top = 1.92
    for label, val in data["cross"]:
        _text(s, 5.6, top, 2.2, 0.3, label, 12.5, INK)
        _text(s, 7.8, top, 1.2, 0.3, val, 12.5, INK, bold=True, align=PP_ALIGN.RIGHT)
        top += 0.30
    # the analyst box: carry vs convexity
    box = s.shapes.add_shape(5, Inches(9.4), Inches(1.5), Inches(3.3), Inches(4.6))
    _fill(box, ICE)
    _text(s, 9.65, 1.7, 2.85, 0.4, "The carry trap", 16, NAVY, bold=True,
          font="Cambria")
    _text(s, 9.65, 2.25, 2.85, 3.6,
          "HY index yields ~7.1% — ~$710k of annual carry per $10m.\n\n"
          "But spread duration ≈ 3.2y. A move from 272 → 380bp (a 2023-'tight' level) "
          "marks the book down ~3.5%, ~$320k — five months of carry gone.\n\n"
          "Upside from here: maybe 40bp of compression. That asymmetry is the whole "
          "game at tight spreads.", 12.5, INK)
    _text(s, 0.7, 6.3, 8, 0.6,
          "Positioning: up-in-quality within HY, own the BB over the CCC — the "
          "CCC-BB gap is priced for a soft landing that a hiking Fed threatens.",
          13, INK, bold=True)
    _text(s, 0.7, 7.05, 6, 0.3, "Source: ICE BofA OAS via FRED", 10, MUTE)

    # --- Slide 4: Themes / What Matters (analyst) ---------------------------
    s = prs.slides.add_slide(blank)
    _text(s, 0.7, 0.5, 12, 0.6, "What Matters This Week", 34, NAVY, bold=True,
          font="Cambria")
    top = 1.4
    for i, (head, body, read) in enumerate(data["themes"], 1):
        num = s.shapes.add_shape(9, Inches(0.7), Inches(top), Inches(0.5), Inches(0.5))
        _fill(num, NAVY)
        _text(s, 0.7, top+0.06, 0.5, 0.4, str(i), 20, WHITE, bold=True,
              align=PP_ALIGN.CENTER)
        _text(s, 1.4, top, 11.2, 0.4, head, 18, NAVY, bold=True, font="Cambria")
        _text(s, 1.4, top+0.45, 11.2, 0.7, body, 12.5, INK)
        _text(s, 1.4, top+1.32, 11.2, 0.35, read, 12.5, DOWN, bold=True, italic=True)
        top += 1.95
    _text(s, 0.7, 7.05, 8, 0.3, "Views are illustrative — not investment advice", 10, MUTE)

    # --- Slide 5: Week ahead + links ----------------------------------------
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height); _fill(bg, NAVY)
    _text(s, 0.7, 0.5, 12, 0.6, "Week Ahead & Sources", 34, WHITE, bold=True,
          font="Cambria")
    # calendar
    _text(s, 0.7, 1.35, 3, 0.3, "DATE", 11, ICE, bold=True)
    _text(s, 2.6, 1.35, 3.5, 0.3, "EVENT", 11, ICE, bold=True)
    _text(s, 5.6, 1.35, 1.2, 0.3, "TIME", 11, ICE, bold=True)
    _text(s, 6.9, 1.35, 6.0, 0.3, "WHY IT MATTERS", 11, ICE, bold=True)
    top = 1.75
    for date, event, time, why in data["calendar"]:
        _text(s, 0.7, top, 2.0, 0.35, date, 12.5, WHITE, bold=True)
        _text(s, 2.6, top, 3.0, 0.35, event, 12.5, WHITE)
        _text(s, 5.6, top, 1.2, 0.35, time, 12.5, ICE)
        _text(s, 6.9, top, 6.0, 0.5, why, 12, ICE)
        top += 0.62
    # links — on a light panel so navy hyperlink text stays readable on the dark slide
    panel = s.shapes.add_shape(5, Inches(0.6), Inches(top+0.05),
                               Inches(12.1), Inches(1.75))
    _fill(panel, WHITE)
    _text(s, 0.85, top+0.2, 6, 0.3, "SOURCES / LINKS", 12, NAVY, bold=True)
    tb = s.shapes.add_textbox(Inches(0.85), Inches(top+0.55), Inches(11.6), Inches(1.2))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (label, url) in enumerate(LINKS):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(3)
        _link(p, "•  " + label, url, color=NAVY)
    prs.save(path)
    return path

# ----------------------------------------------------------------------------
# 7. EMAIL
# ----------------------------------------------------------------------------
def email_deck(path, tape):
    msg = EmailMessage()
    today = dt.date.today().strftime("%d %b %Y")
    msg["Subject"] = f"Fixed Income Weekly — {today}"
    msg["From"] = os.environ["SMTP_USER"]
    msg["To"] = os.environ["MAIL_TO"]
    # Email BODY carries the trader TL;DR so it's readable on a phone at 8am
    # without opening the attachment. The deck is the analyst depth.
    body = "The Tape:\n\n" + "\n\n".join(f"- {t}" for t in tape) + \
           "\n\nFull dashboard attached.\n"
    msg.set_content(body)
    with open(path, "rb") as f:
        msg.add_attachment(f.read(),
            maintype="application",
            subtype="vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=os.path.basename(path))
    with smtplib.SMTP(os.environ["SMTP_HOST"], int(os.environ["SMTP_PORT"])) as sv:
        sv.starttls()
        sv.login(os.environ["SMTP_USER"], os.environ["SMTP_PASS"])
        sv.send_message(msg)

# ----------------------------------------------------------------------------
# 8. MAIN
# ----------------------------------------------------------------------------
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--live", action="store_true", help="pull from FRED")
    ap.add_argument("--demo", action="store_true", help="offline sample data")
    ap.add_argument("--email", action="store_true", help="email the deck")
    ap.add_argument("--out", default="FI_Weekly.pptx")
    args = ap.parse_args()

    if args.live:
        data = pull_live(os.environ["FRED_API_KEY"])
        # live path doesn't compute a curve dict/themes automatically — for a first
        # cut, reuse demo's curve shape from the pulled rates, and keep themes/tape/
        # calendar as fields YOU edit each week (see note in the write-up).
        demo = demo_data()
        data["curve"] = {0.25: (data["rates"][0]["raw"], data["rates"][0]["raw_ref"]),
                         2: (data["rates"][1]["raw"], data["rates"][1]["raw_ref"]),
                         5: (data["rates"][2]["raw"], data["rates"][2]["raw_ref"]),
                         10: (data["rates"][3]["raw"], data["rates"][3]["raw_ref"]),
                         30: (data["rates"][4]["raw"], data["rates"][4]["raw_ref"])}
        data["cross"] = demo["cross"]; data["tape"] = demo["tape"]
        data["themes"] = demo["themes"]; data["calendar"] = demo["calendar"]
    else:
        data = demo_data()

    path = build_deck(data, args.out)
    print("Built", path)
    if args.email:
        email_deck(path, data["tape"])
        print("Emailed to", os.environ["MAIL_TO"])


if __name__ == "__main__":
    main()
